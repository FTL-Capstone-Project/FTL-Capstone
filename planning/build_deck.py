#!/usr/bin/env python3
"""
Generate Orbis_Demo.pptx — Week 6 capstone demo deck for team DOMinion.
Visual-first, ~6 minutes, 3 presenters. Slides carry cue points only;
expansion lives in the speaker notes.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Palette (deep-space / security theme — "Orbis / Orbo")
# ----------------------------------------------------------------------------
BG      = RGBColor(0x0B, 0x10, 0x20)
BG2     = RGBColor(0x0E, 0x16, 0x30)
PANEL   = RGBColor(0x16, 0x1F, 0x3D)
PANEL2  = RGBColor(0x1E, 0x2A, 0x4A)
STROKE  = RGBColor(0x2B, 0x3A, 0x63)
CYAN    = RGBColor(0x22, 0xD3, 0xEE)
SKY     = RGBColor(0x38, 0xBD, 0xF8)
VIOLET  = RGBColor(0x8B, 0x5C, 0xF6)
GREEN   = RGBColor(0x34, 0xD3, 0x99)
AMBER   = RGBColor(0xFB, 0xBF, 0x24)
RED     = RGBColor(0xF8, 0x71, 0x71)
TEXT    = RGBColor(0xE6, 0xEA, 0xF3)
MUTED   = RGBColor(0x93, 0xA3, 0xC4)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x0B, 0x10, 0x20)

FONT = "Arial"  # portable across PowerPoint + Google Slides

EMU_IN = 914400
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s

def _noshadow(sp):
    sp.shadow.inherit = False

def rect(s, x, y, w, h, fill=PANEL, line=STROKE, line_w=1.0, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    _noshadow(shp)
    # soften rounded corner radius
    if rounded:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    return shp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=2, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold, italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for run in para:
            t, size, color, bold = run[0], run[1], run[2], run[3]
            italic = run[4] if len(run) > 4 else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.name = FONT
            r.font.color.rgb = color; r.font.bold = bold; r.font.italic = italic
    return tb

def one(text, size, color, bold=False, italic=False):
    return [(text, size, color, bold, italic)]

def title(s, main, kicker=None):
    # accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.11), Inches(0.62))
    bar.fill.solid(); bar.fill.fore_color.rgb = CYAN; bar.line.fill.background(); _noshadow(bar)
    txt(s, 0.85, 0.44, 11.8, 0.9, [one(main, 30, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        txt(s, 0.86, 1.12, 11.8, 0.35, [one(kicker, 13, CYAN, False)], anchor=MSO_ANCHOR.TOP)

def footer(s, presenter, n):
    txt(s, 0.6, 7.02, 8, 0.35, [one("Orbis  ·  team DOMinion", 10, MUTED, False)])
    txt(s, 6.5, 7.02, 5.3, 0.35, [[("▸ ", 10, CYAN, True), (presenter, 10, MUTED, False),
                                    (f"    {n}/11", 10, STROKE, False)]], align=PP_ALIGN.RIGHT)

def chip(s, x, y, w, label, color, h=0.34, fs=10.5):
    c = rect(s, x, y, w, h, fill=None, line=color, line_w=1.25, rounded=True)
    try: c.adjustments[0] = 0.5
    except Exception: pass
    txt(s, x, y, w, h, [one(label, fs, color, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c

def arrow(s, x, y, w, h=0.32, color=CYAN):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); _noshadow(a)
    try:
        a.adjustments[0] = 0.55  # thinner tail
        a.adjustments[1] = 0.55
    except Exception:
        pass
    return a

def down_arrow(s, x, y, w=0.32, h=0.4, color=CYAN):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); _noshadow(a)
    return a

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

# ============================================================================
# SLIDE 1 — Title
# ============================================================================
s = slide(BG)
# faint orbit rings (decorative)
for rr, col in [(3.2, STROKE), (2.3, PANEL2), (1.5, STROKE)]:
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(SW-4.6-rr/2), Inches(SH/2-rr/2), Inches(rr), Inches(rr))
    o.fill.background(); o.line.color.rgb = col; o.line.width = Pt(1.25); _noshadow(o)
core = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(SW-4.6-0.28), Inches(SH/2-0.28), Inches(0.56), Inches(0.56))
core.fill.solid(); core.fill.fore_color.rgb = CYAN; core.line.fill.background(); _noshadow(core)
dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(SW-4.6+1.5-0.13), Inches(SH/2-1.15-0.13), Inches(0.26), Inches(0.26))
dot.fill.solid(); dot.fill.fore_color.rgb = VIOLET; dot.line.fill.background(); _noshadow(dot)

txt(s, 0.9, 2.0, 8, 0.5, [one("TEAM DOMINION  ·  WEEK 6 CAPSTONE DEMO", 14, CYAN, True)])
txt(s, 0.86, 2.5, 8.2, 1.3, [one("Orbis", 76, WHITE, True)])
txt(s, 0.9, 3.85, 8, 0.6, [one("Know before you click.", 24, TEXT, False, True)])
txt(s, 0.9, 4.7, 8.2, 0.9,
    [one("AI-powered phishing triage — a plain-English danger verdict", 15, MUTED, False),
     one("for anyone who gets a suspicious link.", 15, MUTED, False)])
txt(s, 0.9, 6.2, 9, 0.6,
    [[("Michael Jissa", 14, TEXT, True), ("   ·   ", 14, STROKE, False),
      ("Ozias Tumimana", 14, TEXT, True), ("   ·   ", 14, STROKE, False),
      ("David Gonzalez-Cesar", 14, TEXT, True)]])
notes(s, "ALL THREE presenters step to the podium. Michael opens.\n"
         "Introductions: name + one line each. \n"
         "Michael: 'We're team DOMinion. Our capstone is Orbis — an AI-powered phishing "
         "triage tool. The idea in one line: know before you click.'\n"
         "Keep this to ~20s, then go to the problem.")

# ============================================================================
# SLIDE 2 — Problem
# ============================================================================
s = slide(BG)
title(s, "The Problem", "Phishing is the #1 way organizations get breached")
cards = [
    ("Analysts are buried", "Triaging each report by hand — confirm it's malicious, see what it does, spot a campaign — is slow, repetitive, expert work.", AMBER),
    ("Reporters hear nothing", "You report a suspicious link and wait. No quick answer, no closure, no idea if you did the right thing.", SKY),
    ("The most-targeted have the least help", "Students, individuals, and small companies with no security team are on their own.", RED),
]
cw = 3.85; gap = 0.28; x0 = 0.6; y0 = 1.85
for i, (h, b, col) in enumerate(cards):
    x = x0 + i*(cw+gap)
    rect(s, x, y0, cw, 3.4, fill=PANEL, line=STROKE)
    rect(s, x, y0, 0.12, 3.4, fill=col, line=None)  # left accent
    txt(s, x+0.35, y0+0.35, cw-0.6, 1.0, [one(h, 18, WHITE, True)])
    txt(s, x+0.35, y0+1.35, cw-0.6, 1.9, [one(b, 13.5, MUTED, False)], line_spacing=1.08)
txt(s, 0.6, 5.6, 12.1, 0.9,
    [[("The safe move — ", 15, TEXT, False), ("report it and wait for an analyst", 15, TEXT, True),
      (" — is exactly where everything stalls.", 15, TEXT, False)]],
    anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Michael", 2)
notes(s, "Michael.\n"
         "'Phishing is the number-one entry point for breaches. But triage is a bottleneck.'\n"
         "Walk the three cards left to right (one sentence each):\n"
         " 1) Analysts are buried — every report is slow, manual, specialized work.\n"
         " 2) The person who reported gets no answer — no closure.\n"
         " 3) The people most targeted have the least help — no security team at all.\n"
         "Land the bottom line: the safe move is the thing that stalls. ~40s.")

# ============================================================================
# SLIDE 3 — Solution + roles
# ============================================================================
s = slide(BG)
title(s, "Our Solution — Orbis", "Takes the slow, manual work out of phishing triage")
# top: the one-liner pipeline
steps = ["Submit a suspicious URL", "Detonate in a secure sandbox", "Get a plain-English verdict"]
sub   = ["web chat or forward an email", "urlscan.io — nobody opens the link", "0–100 score + screenshot"]
scol  = [SKY, VIOLET, GREEN]
bw = 3.55; y0 = 1.75; x = 0.6
for i, (st, sb, cc) in enumerate(zip(steps, sub, scol)):
    rect(s, x, y0, bw, 1.35, fill=PANEL, line=STROKE)
    txt(s, x+0.28, y0+0.2, bw-0.5, 0.6, [one(st, 15, WHITE, True)])
    txt(s, x+0.28, y0+0.82, bw-0.5, 0.45, [one(sb, 12, cc, False)])
    if i < 2:
        arrow(s, x+bw+0.05, y0+0.5, 0.55, 0.34, CYAN)
    x += bw + 0.6
# bottom: one product, three roles
txt(s, 0.6, 3.55, 12, 0.4, [[("One product, ", 16, TEXT, False), ("three roles", 16, CYAN, True),
                             ("  — each org sees only its own data.", 16, TEXT, False)]])
roles = [
    ("Individual", "A solo user vetting links for themselves. No IT to ask.", "Lightweight “is this safe?” + my history", SKY),
    ("Organization Member", "Gets suspicious links at work; a bad click risks the whole org.", "Same, plus escalation & closure status", VIOLET),
    ("Security Analyst", "Often the only security person at a small company.", "Full triage dashboard, campaigns, Ask-the-data", GREEN),
]
cw = 3.85; gap = 0.28; y0 = 4.1
for i, (h, b, feat, col) in enumerate(roles):
    x = 0.6 + i*(cw+gap)
    rect(s, x, y0, cw, 2.55, fill=PANEL2, line=STROKE)
    rect(s, x, y0, cw, 0.62, fill=col, line=None)
    txt(s, x, y0, cw, 0.62, [one(h, 15, INK, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+0.3, y0+0.85, cw-0.55, 1.0, [one(b, 12.5, MUTED, False)], line_spacing=1.05)
    txt(s, x+0.3, y0+1.85, cw-0.55, 0.6, [[("Gets:  ", 12, col, True), (feat, 12, TEXT, False)]], line_spacing=1.05)
footer(s, "Michael", 3)
notes(s, "Michael.\n"
         "'Orbis is what happens when you automate the slow part.' Walk the top pipeline: "
         "submit a URL (web or email) -> we detonate it in urlscan.io's sandbox so nobody opens it "
         "-> Orbis returns a plain-English verdict, a 0-100 score, and a screenshot.\n"
         "Then the key framing: ONE product serves THREE roles, and each org sees only its own data. "
         "Individual = lightweight self-check; Org Member = same + closure status; Analyst = full triage dashboard.\n"
         "Hand off to David after this slide. ~45s.")

# ============================================================================
# SLIDE 4 — Planned features
# ============================================================================
s = slide(BG)
title(s, "What We're Building — MVP", "Scoped for the sprint; stretch items on the roadmap")
feats = [
    ("Check-a-link verdict", "Submit a URL, sandbox scan, AI score + screenshot", "AI", GREEN),
    ("Email forwarding", "Forward to the Orbo inbox — no UI, same flow", "core", CYAN),
    ("Role-tailored Reports", "Individual / Member / Analyst — three views", "core", CYAN),
    ("Auto-escalation + closure", "Member submissions route to the analyst; notify on verdict", "core", CYAN),
    ("Analyst triage dashboard", "Org-wide history, stats, campaign grouping", "core", CYAN),
    ("Ask Orbo (ask-the-data)", "Plain-English question → chart, no SQL", "AI", GREEN),
    ("Managed auth + orgs (Clerk)", "Social login & org invites out of the box", "core", CYAN),
    ("Known-bad blacklist check", "Google Safe Browsing signal per URL", "core", CYAN),
]
cw = 5.9; ch = 1.02; gx = 0.35; gy = 0.28; x0 = 0.6; y0 = 1.75
for i, (h, b, tag, col) in enumerate(feats):
    r, c = divmod(i, 2)
    x = x0 + c*(cw+gx); y = y0 + r*(ch+gy)
    rect(s, x, y, cw, ch, fill=PANEL, line=STROKE)
    rect(s, x, y, 0.1, ch, fill=col, line=None)
    txt(s, x+0.3, y+0.14, cw-1.4, 0.4, [one(h, 14.5, WHITE, True)])
    txt(s, x+0.3, y+0.56, cw-1.4, 0.4, [one(b, 11.5, MUTED, False)])
    chip(s, x+cw-1.0, y+0.32, 0.8, tag, col, h=0.3, fs=9.5)
txt(s, 0.6, 6.75, 12, 0.4,
    [[("Stretch (roadmap):  ", 12.5, VIOLET, True),
      ("browser extension  ·  enterprise SSO/SAML  ·  email/SMS alerts  ·  more Ask-Orbo charts", 12.5, MUTED, False)]])
footer(s, "David", 4)
notes(s, "David.\n"
         "'Here's what ships in the MVP.' Don't read all eight — group them:\n"
         " - The spine: check-a-link with an AI verdict (green = AI feature).\n"
         " - Two ways in: web chat AND forward-an-email (same backend flow).\n"
         " - The two-sided part: role-tailored reports, auto-escalation to the analyst, closure notifications.\n"
         " - The analyst power tools: triage dashboard + Ask Orbo (our 2nd AI feature).\n"
         " - Bought-not-built: Clerk auth/orgs, Google Safe Browsing blacklist.\n"
         "One line on stretch, then move to the wireframes. ~40s.")

# ============================================================================
# SLIDE 5 — Wireframes / workflow
# ============================================================================
s = slide(BG)
title(s, "Wireframes — the Check-a-Link Flow", "Designed in Figma · full flow drawn · live demo to follow")
flow = [
    ("Home", "Paste a link\n(or forward email)", SKY),
    ("Checking", "Sandbox runs\nloading state", AMBER),
    ("Result", "Score · screenshot\nplain-English verdict", GREEN),
    ("Reports", "My history &\nclosure status", VIOLET),
]
bw = 2.7; y0 = 2.2; x = 0.7
for i, (h, b, col) in enumerate(flow):
    rect(s, x, y0, bw, 1.9, fill=PANEL, line=STROKE)
    rect(s, x, y0, bw, 0.55, fill=col, line=None)
    txt(s, x, y0, bw, 0.55, [one(h, 15, INK, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for j, ln in enumerate(b.split("\n")):
        txt(s, x+0.2, y0+0.75+j*0.34, bw-0.4, 0.34, [one(ln, 12.5, TEXT, False)], align=PP_ALIGN.CENTER)
    if i < 3:
        arrow(s, x+bw+0.02, y0+0.78, 0.5, 0.34, CYAN)
    x += bw + 0.52
# verdict states row
txt(s, 0.7, 4.55, 12, 0.4, [one("One verdict card — three data-driven states:", 14, TEXT, True)])
for i, (lbl, col) in enumerate([("SAFE", GREEN), ("SUSPICIOUS", AMBER), ("DANGEROUS", RED)]):
    chip(s, 0.7 + i*2.0, 5.1, 1.8, lbl, col, h=0.42, fs=12)
txt(s, 6.6, 5.05, 6.1, 1.3,
    [one("Also drawn: landing, login/register, team setup,", 12.5, MUTED, False),
     one("analyst dashboard, Ask Orbo (6 chart variants),", 12.5, MUTED, False),
     one("3 report variants + detail modal, invalid-input state.", 12.5, MUTED, False)],
    line_spacing=1.1)
footer(s, "David", 5)
notes(s, "David — this is the hand-off into the live wireframe walkthrough (NOT in the slides).\n"
         "'Our core flow is four screens.' Trace it: Home (paste or forward) -> Checking "
         "(sandbox running) -> Result (score, screenshot, plain-English verdict) -> Reports (history + closure).\n"
         "Point out the single verdict card renders Safe / Suspicious / Dangerous from the data.\n"
         "Mention we drew far more than the required 3 screens, then say: 'Let me show you the actual "
         "wireframes' and switch to Figma/PDF. Come back to the deck for the data model. ~40s + demo.")

# ============================================================================
# SLIDE 6 — Data model (visual, two layers)
# ============================================================================
s = slide(BG)
title(s, "Data Model — Two Layers", "Shared threat intel  ·  private per-org data")
# Clerk mirror badge (top-right)
rect(s, 10.05, 1.62, 2.65, 0.95, fill=PANEL2, line=VIOLET, line_w=1.25)
txt(s, 10.05, 1.7, 2.65, 0.4, [one("CLERK (auth)", 12, VIOLET, True)], align=PP_ALIGN.CENTER)
txt(s, 10.1, 2.08, 2.55, 0.45, [one("mirrors → users · organizations", 10.5, MUTED, False)], align=PP_ALIGN.CENTER)

# GLOBAL layer
gy = 1.62
rect(s, 0.6, gy, 9.2, 2.15, fill=PANEL, line=CYAN, line_w=1.5)
txt(s, 0.85, gy+0.12, 9, 0.4, [[("GLOBAL — shared threat intel", 14, CYAN, True),
                                (" · scanned once, reused by everyone", 12, MUTED, False)]])
# indicator card
rect(s, 0.9, gy+0.62, 3.3, 1.35, fill=PANEL2, line=CYAN)
txt(s, 1.05, gy+0.72, 3.0, 0.35, [one("indicators", 13.5, WHITE, True)])
txt(s, 1.05, gy+1.12, 3.05, 0.85,
    [one("canonical_key (unique)", 10.5, CYAN, False),
     one("ai_score · ai_verdict · screenshot", 10.5, MUTED, False),
     one("blacklist_hit · report_count", 10.5, MUTED, False)], line_spacing=1.05)
txt(s, 4.5, gy+0.72, 5.1, 1.2,
    [one("One row per unique URL.", 12.5, TEXT, True),
     one("The objective facts about a link —", 11.5, MUTED, False),
     one("is it malicious, the AI reasoning, the", 11.5, MUTED, False),
     one("screenshot — are the same for everyone.", 11.5, MUTED, False)], line_spacing=1.05)

down_arrow(s, 5.0, 3.85, 0.34, 0.34, STROKE)

# PER-ORG layer
py = 4.35
rect(s, 0.6, py, 12.1, 2.35, fill=PANEL, line=VIOLET, line_w=1.5)
txt(s, 0.85, py+0.12, 11.5, 0.4, [[("PER-ORG — private", 14, VIOLET, True),
                                   (" · walled off by org_id (story #12)", 12, MUTED, False)]])
tables = [
    ("submissions", "who reported it\nraw_url · source", SKY),
    ("org_reviews", "analyst's verdict\nhuman_score · status", GREEN),
    ("campaigns", "cluster of threats\nshared_signal", AMBER),
    ("notifications", "closure alerts\nto the reporter", VIOLET),
]
cw = 2.82; x = 0.9
for (h, b, col) in tables:
    rect(s, x, py+0.62, cw, 1.55, fill=PANEL2, line=col)
    rect(s, x, py+0.62, cw, 0.45, fill=col, line=None)
    txt(s, x, py+0.62, cw, 0.45, [one(h, 12.5, INK, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for j, ln in enumerate(b.split("\n")):
        txt(s, x+0.15, py+1.2+j*0.34, cw-0.3, 0.34, [one(ln, 10.5, MUTED if j else TEXT, False)], align=PP_ALIGN.CENTER)
    x += cw + 0.2
footer(s, "Ozias", 6)
notes(s, "Ozias.\n"
         "'The one design decision that makes Orbis work is splitting the data into two layers.'\n"
         "TOP (global indicators): the objective facts about a URL — malicious or not, the AI verdict, "
         "the screenshot. That's threat INTEL, not private data, so we scan a URL ONCE and reuse it for "
         "everyone. That powers 'Orbo has seen this before, reported N times.'\n"
         "BOTTOM (per-org): who reported it, the analyst's authoritative verdict, campaigns, notifications — "
         "all walled off by org_id so one company can never see another's activity.\n"
         "Right badge: Clerk owns identity; users & organizations are just mirror rows. ~50s. Go to flow.")

# ============================================================================
# SLIDE 7 — Data flow
# ============================================================================
s = slide(BG)
title(s, "How the Data Flows", "From a pasted link to a closed investigation")
# main horizontal pipeline
nodes = [
    ("Submit", "web or email", SKY),
    ("Dedup", "canonical_key\ncollapse duplicates", CYAN),
    ("Scan + Blacklist", "urlscan.io +\nSafe Browsing", AMBER),
    ("AI Verdict", "score + plain\nEnglish (Claude)", GREEN),
]
bw = 2.72; y0 = 1.95; x = 0.6
for i, (h, b, col) in enumerate(nodes):
    rect(s, x, y0, bw, 1.5, fill=PANEL, line=col, line_w=1.4)
    txt(s, x+0.2, y0+0.16, bw-0.4, 0.4, [one(h, 14.5, WHITE, True)], align=PP_ALIGN.CENTER)
    for j, ln in enumerate(b.split("\n")):
        txt(s, x+0.2, y0+0.66+j*0.32, bw-0.4, 0.32, [one(ln, 11, MUTED, False)], align=PP_ALIGN.CENTER)
    if i < 3:
        arrow(s, x+bw+0.01, y0+0.58, 0.42, 0.32, CYAN)
    x += bw + 0.44

# fork down: two paths
txt(s, 0.6, 3.75, 12, 0.35, [one("Then it forks on who submitted:", 13.5, TEXT, True)])
# individual path
rect(s, 0.6, 4.25, 5.85, 2.1, fill=PANEL2, line=SKY)
txt(s, 0.85, 4.4, 5.4, 0.4, [one("Individual  (org_id = NULL)", 13.5, SKY, True)])
txt(s, 0.85, 4.9, 5.4, 1.3,
    [one("• Gets the AI verdict instantly.", 12.5, TEXT, False),
     one("• No analyst, nothing escalates.", 12.5, MUTED, False),
     one("• Scoped to their own user_id.", 12.5, MUTED, False)], line_spacing=1.2)
# org member path
rect(s, 6.85, 4.25, 5.85, 2.1, fill=PANEL2, line=GREEN)
txt(s, 7.1, 4.4, 5.4, 0.4, [one("Org Member  →  Analyst", 13.5, GREEN, True)])
txt(s, 7.1, 4.9, 5.5, 1.5,
    [[("• Auto-escalated ", 12.5, TEXT, False), ("(escalated = true)", 11, MUTED, False)],
     one("• Analyst writes authoritative verdict", 12.5, TEXT, False),
     [("• ", 12.5, TEXT, False), ("Notification", 12.5, GREEN, True), (" → reporter gets closure", 12.5, TEXT, False)]],
    line_spacing=1.2)
footer(s, "Ozias", 7)
notes(s, "Ozias.\n"
         "Trace the top pipeline left to right: a user submits (web or email) -> we compute a "
         "canonical_key to collapse duplicate/tracking-junk URLs into one indicator -> if it's new we scan it "
         "in urlscan AND check Google Safe Browsing -> Claude turns the evidence into a score + plain-English verdict.\n"
         "Then it FORKS on who submitted:\n"
         " - Individual (no org): gets the verdict, done. Nothing escalates.\n"
         " - Org member: auto-escalated to the analyst; the analyst records the authoritative verdict; "
         "a notification fires back so the reporter gets closure (story #7).\n"
         "Emphasize: the dedup is what turns 20 reports into 1 investigation. ~50s.")

# ============================================================================
# SLIDE 8 — API endpoints
# ============================================================================
s = slide(BG)
title(s, "API — the Endpoints We Build", "12 first-party routes (auth/orgs handled by Clerk)")
def ep(s, x, y, w, verb, path, desc, vcol):
    rect(s, x, y, w, 0.62, fill=PANEL, line=STROKE)
    chip(s, x+0.12, y+0.14, 0.85, verb, vcol, h=0.34, fs=9.5)
    txt(s, x+1.08, y+0.08, w-1.2, 0.3, [one(path, 11.5, WHITE, True)])
    txt(s, x+1.08, y+0.34, w-1.2, 0.28, [one(desc, 10, MUTED, False)])
POST=RGBColor(0x34,0xD3,0x99); GET=RGBColor(0x38,0xBD,0xF8); PATCH=RGBColor(0xFB,0xBF,0x24)
left = [
    ("POST","/api/submissions","Submit a URL for analysis", POST),
    ("GET","/api/indicators/:id","Poll the verdict (global + my review)", GET),
    ("GET","/api/history?mine=1","My reported links", GET),
    ("GET","/api/history","Org-wide history + stats (analyst)", GET),
    ("GET","/api/search?q=","Keyword search in my org (analyst)", GET),
    ("PATCH","/api/indicators/:id/review","Analyst records authoritative verdict", PATCH),
]
right = [
    ("GET","/api/campaigns","My org's campaigns (triage queue)", GET),
    ("POST","/api/nlp-query","★ Ask Orbo: English → chart", POST),
    ("GET","/api/notifications","My closure alerts", GET),
    ("POST","/api/webhooks/clerk","Sync user/org mirror rows", POST),
    ("POST","/api/webhooks/inbound-email","Orbo inbox → submission", POST),
    ("GET","/api/campaigns/:id","Campaign detail (grouped)", GET),
]
y = 1.8
for e in left:  ep(s, 0.6, y, 5.95, *e); y += 0.72
y = 1.8
for e in right: ep(s, 6.75, y, 5.95, *e); y += 0.72
txt(s, 0.6, 6.55, 12.1, 0.7,
    [[("One middleware ", 12.5, CYAN, True),
      ("verifies the Clerk session and enforces role + org_id on every protected route — ", 12.5, TEXT, False),
      ("so no org can ever read another's data.", 12.5, TEXT, True)]], line_spacing=1.1)
footer(s, "Ozias", 8)
notes(s, "Ozias.\n"
         "Don't read all twelve — call out the shape: full CRUD, comfortably past the 5-endpoint bar.\n"
         "Highlight four: POST /submissions (the entry point), GET /indicators/:id (poll for the verdict — "
         "note we merge the GLOBAL indicator with the caller's private review), PATCH .../review (the "
         "analyst's authoritative verdict — the two-sided part), and POST /nlp-query (the star — Ask Orbo).\n"
         "Note the two webhooks: Clerk sync + inbound email.\n"
         "Land the bottom line: ONE middleware enforces role + org on every route — that's how isolation is real, "
         "not just a promise. Hand back to David/Michael. ~50s.")

# ============================================================================
# SLIDE 9 — AI features
# ============================================================================
s = slide(BG)
title(s, "Two AI Features", "Rubric asks for one — Orbis has two")
# Feature A
rect(s, 0.6, 1.8, 5.95, 4.5, fill=PANEL, line=GREEN, line_w=1.5)
rect(s, 0.6, 1.8, 5.95, 0.7, fill=GREEN, line=None)
txt(s, 0.6, 1.8, 5.95, 0.7, [one("A · Plain-English Danger Verdict", 15, INK, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.9, 2.75, 5.4, 3.4,
    [[("Turns ", 13, MUTED, False), ("raw sandbox evidence", 13, TEXT, True), (" into a", 13, MUTED, False)],
     one("verdict + 0–100 score anyone can act on.", 13, MUTED, False),
     one("", 6, MUTED, False),
     [("In:  ", 12.5, GREEN, True), ("scan evidence + blacklist + context", 12, TEXT, False)],
     [("Out:  ", 12.5, GREEN, True), ("structured JSON (always valid)", 12, TEXT, False)],
     one("", 6, MUTED, False),
     [("Safety floor:  ", 12.5, GREEN, True), ("a known-bad URL can", 12, TEXT, False)],
     one("never be reported as “safe.”", 12, TEXT, False)], line_spacing=1.12)
# Feature B
rect(s, 6.75, 1.8, 5.95, 4.5, fill=PANEL, line=SKY, line_w=1.5)
rect(s, 6.75, 1.8, 5.95, 0.7, fill=SKY, line=None)
txt(s, 6.75, 1.8, 5.95, 0.7, [one("B · Ask Orbo — Ask-the-Data", 15, INK, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 7.05, 2.75, 5.4, 3.4,
    [[("An analyst asks in ", 13, MUTED, False), ("plain English", 13, TEXT, True), (",", 13, MUTED, False)],
     one("gets a chart back — no SQL.", 13, MUTED, False),
     one("", 6, MUTED, False),
     [("In:  ", 12.5, SKY, True), ("question + allowed field schema", 12, TEXT, False)],
     [("Out:  ", 12.5, SKY, True), ("validated filter + chartSpec", 12, TEXT, False)],
     one("", 6, MUTED, False),
     [("Safety:  ", 12.5, SKY, True), ("whitelisted filter, not raw SQL", 12, TEXT, False)],
     one("— no AI-driven injection path.", 12, TEXT, False)], line_spacing=1.12)
footer(s, "David", 9)
notes(s, "David.\n"
         "'The rubric needs one AI feature — we have two, and both are hardened.'\n"
         "A) The verdict: Claude turns sandbox evidence into a human-readable score. Two safety moves — "
         "structured outputs (always valid JSON) and a deterministic FLOOR so a blacklisted URL can never "
         "come back 'safe,' no matter what the model says.\n"
         "B) Ask Orbo: analyst types a question, gets a chart. The model emits a WHITELISTED FILTER OBJECT, "
         "not SQL — so there's no injection path; the backend runs a parameterized query.\n"
         "This sets up the next slide (challenge). ~45s.")

# ============================================================================
# SLIDE 10 — Most interesting & challenging
# ============================================================================
s = slide(BG)
title(s, "Most Interesting & Challenging", "The dedup that makes everything else work")
rect(s, 0.6, 1.8, 12.1, 1.55, fill=PANEL2, line=CYAN, line_w=1.5)
txt(s, 0.9, 1.95, 11.6, 0.5, [one("The problem: phishers give every victim a “unique” link.", 15, WHITE, True)])
txt(s, 0.9, 2.5, 11.6, 0.8,
    [[(".../verify?id=", 13, MUTED, False), ("david", 13, RED, True), ("&ref=email123", 13, MUTED, False),
      ("      vs      ", 13, STROKE, False),
      (".../verify?id=", 13, MUTED, False), ("maria", 13, RED, True), ("&ref=email456", 13, MUTED, False),
      ("   →   same attack.", 13, GREEN, True)]])
# why it matters + how
rect(s, 0.6, 3.6, 5.95, 2.75, fill=PANEL, line=STROKE)
txt(s, 0.85, 3.75, 5.5, 0.4, [one("Why it's the crux", 14, CYAN, True)])
txt(s, 0.85, 4.25, 5.5, 2.0,
    [one("• 20 reports → 1 investigation", 12.5, TEXT, False),
     one("• powers “seen before, reported N×”", 12.5, TEXT, False),
     one("• a URL is never scanned twice", 12.5, TEXT, False),
     one("• it's what defines a campaign", 12.5, TEXT, False)], line_spacing=1.35)
rect(s, 6.75, 3.6, 5.95, 2.75, fill=PANEL, line=STROKE)
txt(s, 7.0, 3.75, 5.5, 0.4, [one("Why it's hard", 14, AMBER, True)])
txt(s, 7.0, 4.25, 5.5, 2.0,
    [one("• strip tracking junk, keep meaning", 12.5, TEXT, False),
     one("• denylist is heuristic — must tune", 12.5, TEXT, False),
     one("• shorteners need the resolved URL", 12.5, TEXT, False),
     one("• too aggressive = wrong merges", 12.5, TEXT, False)], line_spacing=1.35)
footer(s, "Michael", 10)
notes(s, "Michael.\n"
         "'The feature we're most excited about — and most worried about — is the dedup key.'\n"
         "Set it up with the example: phishers add per-victim junk so every target gets a different-looking "
         "URL, but it's the same attack. Our canonical_key normalizes them (lowercase host, drop fragments, "
         "strip tracking params, sort the rest) so they collapse into ONE indicator.\n"
         "Why it's the crux: it turns 20 reports into 1 investigation, powers 'seen this before,' avoids "
         "re-scanning, and is literally how we define a campaign.\n"
         "Why it's hard: knowing which params are junk vs. meaningful is heuristic; too aggressive and we "
         "merge unrelated pages; shorteners need the resolved URL (roadmap). ~50s.")

# ============================================================================
# SLIDE 11 — Setup + conclusion
# ============================================================================
s = slide(BG)
title(s, "Project Setup & What's Next", None)
# stack row
txt(s, 0.6, 1.55, 12, 0.4, [one("Stack", 14, CYAN, True)])
stack = ["React", "Node / Express", "PostgreSQL", "Clerk (auth)", "Render (deploy)", "Claude API"]
x = 0.6
for st in stack:
    w = 0.35 + len(st)*0.115
    chip(s, x, 2.0, w, st, TEXT, h=0.42, fs=12)
    x += w + 0.25
# setup + status
rect(s, 0.6, 2.8, 5.95, 2.0, fill=PANEL, line=STROKE)
txt(s, 0.85, 2.95, 5.5, 0.4, [one("Set up & ready", 14, GREEN, True)])
txt(s, 0.85, 3.45, 5.5, 1.4,
    [one("• Repo + GitHub Project board", 12.5, TEXT, False),
     one("• Issues & milestones per sprint", 12.5, TEXT, False),
     one("• Deploy on free tiers (early Sprint 1)", 12.5, TEXT, False),
     one("• Seed script for a realistic demo", 12.5, TEXT, False)], line_spacing=1.3)
rect(s, 6.75, 2.8, 5.95, 2.0, fill=PANEL, line=STROKE)
txt(s, 7.0, 2.95, 5.5, 0.4, [one("Roadmap (stretch)", 14, VIOLET, True)])
txt(s, 7.0, 3.45, 5.5, 1.4,
    [one("• Browser extension", 12.5, TEXT, False),
     one("• Enterprise SSO / SAML", 12.5, TEXT, False),
     one("• Email / SMS alerts", 12.5, TEXT, False),
     one("• More Ask-Orbo chart types", 12.5, TEXT, False)], line_spacing=1.3)
# closing line
rect(s, 0.6, 5.15, 12.1, 1.35, fill=PANEL2, line=CYAN, line_w=1.5)
txt(s, 0.6, 5.25, 12.1, 0.55, [one("Orbis turns a scary link into a clear answer — in seconds.", 18, WHITE, True)],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.6, 5.85, 12.1, 0.5,
    [[("Thank you.  ", 14, CYAN, True), ("Questions?", 14, TEXT, True),
      ("      Michael · Ozias · David — team DOMinion", 12.5, MUTED, False)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, "All", 11)
notes(s, "Michael closes; all three field Q&A.\n"
         "Setup: repo + GitHub Project board with issues/milestones (submission requirement), deploy on free "
         "tiers early in Sprint 1, seed script so the demo looks populated.\n"
         "One line on the roadmap (extension, SSO, alerts, more charts).\n"
         "Closing line — say it with conviction: 'Orbis turns a scary link into a clear answer, in seconds.' "
         "Then invite questions and manage Q&A professionally: whoever owns the area answers; keep answers short.")

# ----------------------------------------------------------------------------
out = "/Users/mjissa/codepath/FTL-Capstone/planning/Orbis_Demo.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
