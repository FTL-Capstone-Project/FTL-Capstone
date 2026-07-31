#!/usr/bin/env python3
"""
Build the Orbis FINAL intern-demo PowerPoint (6-8 min, 3 presenters).

Run:  python3 scripts/build_final_deck.py
Out:  Orbis_Capstone_Final.pptx   (repo root)

WHY THIS DECK IS SHAPED THIS WAY
  Instructor steer: don't put every feature on a slide, narrate it, THEN show the
  same feature again in the demo. Features live INSIDE the live demo. So:
    * BEFORE the demo  -> framing only (problem, how it works, who it's for). No feature tour.
    * THE demo         -> every feature is shown live, driven in the real app.
    * AFTER the demo   -> the ideas you CAN'T click (the deterministic core), the honest
                          roadmap, and a conclusion that ties the bow.

  This mirrors the attached FINAL deck's dark theme, but pulls its colors straight from the
  app's own token file (client/src/theme/tokens.css, dark mode) so the deck and the product
  match exactly. Mascot + logo are the real app assets.

Design method: ASSERTION-EVIDENCE. Each headline is a full-sentence takeaway; the body is
sparse; the DETAIL + per-person script lives in the speaker notes (and in the companion doc
planning/DEMO_SCRIPT.md).

Presentation-only tooling. Touches no application code. Re-run any time after editing CONTENT.
"""

import os
import struct

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --------------------------------------------------------------------------- #
# Paths — real app assets so the deck matches the product
# --------------------------------------------------------------------------- #
REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(REPO, "client", "src", "assets")
ORBO = os.path.join(ASSETS, "orbo")
OUT = os.path.join(REPO, "Orbis_Capstone_Final.pptx")

LOGO_DARK = os.path.join(ASSETS, "orbis-logo-dark.png")   # light-blue logo, for dark slides
M_WAVE    = os.path.join(ORBO, "orbo-wave.png")
M_HAPPY   = os.path.join(ORBO, "orbo-happy.png")
M_THINK   = os.path.join(ORBO, "orbo-thinking.png")
M_CAUTION = os.path.join(ORBO, "orbo-caution.png")
M_DANGER  = os.path.join(ORBO, "orbo-danger.png")
M_SAFE    = os.path.join(ORBO, "orbo-safe.png")

# --------------------------------------------------------------------------- #
# Palette — lifted from client/src/theme/tokens.css  [data-theme="dark"]
# --------------------------------------------------------------------------- #
def C(hexstr):
    return RGBColor.from_string(hexstr)

# base (dark theme)
CANVAS  = C("0B1220")   # --canvas  : page base
SURFACE = C("131C2B")   # --surface : cards, panels
PANEL   = C("0F1830")   # slightly-lifted hero side panel (for depth)
BORDER  = C("24304A")   # --border  : card edges / hairlines
TEXT    = C("E4E8EF")   # --text    : body (~87% white)
DIM     = C("93A0B5")   # --text-dim: secondary / labels
HAZE    = C("9FB0C9")   # softer secondary on hero
WHITE   = C("FFFFFF")

# accents
ACCENT  = C("38BDF8")   # bright cyan kicker/label (matches the PDF's cyan headers)
RING    = C("21C7E6")   # --ring : Orbo's orbital cyan
PRIMARY = C("3B82F6")   # --primary (dark) : links / primary blue

# semantic verdict bands (dark theme)
SAFE    = C("34D07F")   # --safe
REVIEW  = C("E0A93A")   # --review
DANGER  = C("F26B72")   # --danger

# role accents (match the PDF's three-experiences columns)
R_INDIV = C("6366F1")   # individuals — indigo
R_MEMBR = C("21C7E6")   # org members — cyan
R_ANALY = C("34D07F")   # analysts    — green

# presenter accents (lightened for legibility on navy)
MICHAEL = C("34D399")   # emerald
DAVID   = C("60A5FA")   # blue
OZIAS   = C("A78BFA")   # violet

# Fonts (present on macOS; matches the app + the existing pod-sync deck)
DISPLAY = "Avenir Next"
BODY    = "Avenir Next"
MONO    = "Menlo"

# 16:9 canvas
SW = Inches(13.333)
SH = Inches(7.5)
MARGIN = Inches(0.85)
CW = SW - 2 * MARGIN

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #
def png_size(path):
    try:
        with open(path, "rb") as f:
            head = f.read(24)
        if len(head) == 24 and head[:8] == b"\x89PNG\r\n\x1a\n":
            w, h = struct.unpack(">II", head[16:24])
            return w, h
    except OSError:
        pass
    return None


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def no_autofit(tf):
    """Disable PowerPoint auto-shrink so our sizes render as authored."""
    el = tf._txBody
    bodyPr = el.find(qn("a:bodyPr"))
    if bodyPr is None:
        return
    for tag in ("a:normAutofit", "a:spAutoFit"):
        e = bodyPr.find(qn(tag))
        if e is not None:
            bodyPr.remove(e)
    bodyPr.append(el.makeelement(qn("a:noAutofit"), {}))


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    no_autofit(tf)
    return tf


def run(p, text, size=16, color=TEXT, bold=False, italic=False,
        font=BODY, spacing=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    if spacing is not None:
        r.font._rPr.set("spc", str(int(spacing * 100)))
    return r


def para(tf, first=False, align=PP_ALIGN.LEFT, before=0, after=6, line=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.line_spacing = line
    return p


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def image_fit(slide, path, bx, by, bw, bh):
    """Place an image contained within a box, preserving aspect ratio, centered."""
    size = png_size(path)
    if not os.path.exists(path) or size is None:
        return None
    iw, ih = size
    scale = min(bw / iw, bh / ih)
    w, h = Emu(int(iw * scale)), Emu(int(ih * scale))
    x = Emu(int(bx + (bw - w) / 2))
    y = Emu(int(by + (bh - h) / 2))
    return slide.shapes.add_picture(path, x, y, width=w, height=h)


# --------------------------------------------------------------------------- #
# Shared furniture
# --------------------------------------------------------------------------- #
def dark_base(slide, right_panel=True):
    """Every content slide: navy canvas, left cyan spine, optional side panel."""
    set_bg(slide, CANVAS)
    rect(slide, 0, 0, Inches(0.16), SH, fill=RING)  # left orbital-cyan spine
    if right_panel:
        rect(slide, SW - Inches(4.3), 0, Inches(4.3), SH, fill=PANEL)
        rect(slide, SW - Inches(4.3), 0, Pt(1.25), SH, fill=BORDER)


def kicker_head(slide, kicker, headline, hsize=30, accent=ACCENT, top=Inches(0.72),
                width=None):
    """Cyan kicker label + big white assertion headline + accent underline."""
    w = width or (CW - Inches(3.6))
    tf = textbox(slide, MARGIN, top, w, Inches(0.35))
    run(para(tf, first=True, after=0), kicker.upper(), size=13, color=accent,
        bold=True, spacing=2.4)
    htf = textbox(slide, MARGIN, top + Inches(0.45), w, Inches(1.5))
    hp = para(htf, first=True, after=0, line=1.03)
    run(hp, headline, size=hsize, color=WHITE, bold=True, font=DISPLAY)
    rule_y = top + (Inches(1.45) if len(headline) < 54 else Inches(1.95))
    rect(slide, MARGIN, rule_y, Inches(1.0), Pt(3), fill=accent)
    return rule_y + Inches(0.3)


def footer(slide, n):
    tf = textbox(slide, MARGIN, Inches(7.04), Inches(7), Inches(0.3),
                 anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, first=True, after=0), "Orbis  ·  Team DOMinion", size=9, color=DIM,
        spacing=0.5)
    tf2 = textbox(slide, SW - MARGIN - Inches(1.2), Inches(7.04), Inches(1.2),
                  Inches(0.3), anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf2, first=True, align=PP_ALIGN.RIGHT, after=0), f"{n:02d} / 09", size=9,
        color=DIM, spacing=0.5)


def presenter_pill(slide, name, accent):
    """Small pill, top-right of the side panel, showing who drives this slide."""
    w, h = Inches(2.9), Inches(0.44)
    x = SW - MARGIN - w
    rect(slide, x, Inches(0.72), w, h, fill=SURFACE, line=BORDER, line_w=1.25,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, x + Inches(0.18), Inches(0.72) + Inches(0.14), Inches(0.16),
         Inches(0.16), fill=accent, shape=MSO_SHAPE.OVAL)
    tf = textbox(slide, x + Inches(0.46), Inches(0.72), w - Inches(0.56), h,
                 anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True, after=0)
    run(p, name, size=12, color=TEXT, bold=True)
    run(p, "   presents", size=12, color=DIM)


# --------------------------------------------------------------------------- #
# Hero (full-bleed navy) slides
# --------------------------------------------------------------------------- #
def hero_bg(slide, accent=RING):
    set_bg(slide, CANVAS)
    rect(slide, 0, 0, Inches(0.22), SH, fill=accent)
    rect(slide, SW - Inches(4.9), 0, Inches(4.9), SH, fill=PANEL)
    rect(slide, SW - Inches(4.9), 0, Pt(1.5), SH, fill=BORDER)


# --------------------------------------------------------------------------- #
# CONTENT
# --------------------------------------------------------------------------- #
def build():
    n = 0

    # ======================================================================= #
    # 1. TITLE (hero)  — Michael opens
    # ======================================================================= #
    n += 1
    s = add_slide()
    hero_bg(s)
    image_fit(s, LOGO_DARK, Inches(0.9), Inches(0.85), Inches(3.3), Inches(1.0))
    tf = textbox(s, Inches(0.92), Inches(2.35), Inches(8.0), Inches(0.35))
    run(para(tf, first=True, after=0), "CODEPATH SITE CAPSTONE  ·  2026", size=13,
        color=ACCENT, bold=True, spacing=2.4)
    tf2 = textbox(s, Inches(0.9), Inches(2.9), Inches(8.2), Inches(2.0))
    run(para(tf2, first=True, after=0, line=1.02),
        "Stop phishing before\nit ever reaches you.", size=52, color=WHITE, bold=True,
        font=DISPLAY)
    rect(s, Inches(0.95), Inches(4.75), Inches(2.0), Pt(3), fill=RING)
    tf3 = textbox(s, Inches(0.95), Inches(5.05), Inches(7.6), Inches(1.0))
    run(para(tf3, first=True, after=0, line=1.2),
        "AI-assisted phishing triage — expert-grade for analysts, "
        "one-tap simple for everyone else.", size=17, color=HAZE)
    tf4 = textbox(s, Inches(0.95), Inches(6.35), Inches(9.2), Inches(0.5))
    p4 = para(tf4, first=True, after=0)
    run(p4, "Team DOMinion", size=14, color=WHITE, bold=True)
    run(p4, "     Michael Jissa  ·  Ozias Tumimana  ·  David Gonzalez-Cesar",
        size=13, color=HAZE)
    image_fit(s, M_WAVE, SW - Inches(4.55), Inches(3.15), Inches(3.7), Inches(3.7))
    footer(s, n)
    notes(s, """
MICHAEL opens (~15s). Don't read the slide. Set the room in one breath:
"We're Team DOMinion, and Orbis is AI-assisted phishing triage — expert-grade for the analysts
who triage threats, one-tap simple for everyone else. In the next few minutes we'll show you the
problem, then take you straight into the live product."
Then advance. Keep the agenda to ONE sentence — the demo is the star, not the slides.
""")

    # ======================================================================= #
    # 2. THE PROBLEM (dark)  — Michael
    # ======================================================================= #
    n += 1
    s = add_slide()
    dark_base(s)
    presenter_pill(s, "Michael", MICHAEL)
    y = kicker_head(s, "The problem",
                    "Phishing is a human problem\nsolved only with expert tools.",
                    hsize=29)
    # three stat tiles
    stats = [("$4.9B", "reported phishing losses\nin a single year", DANGER),
             ("~90%", "of breaches begin with\na phishing email", REVIEW),
             ("1 click", "is all it takes to detonate\non your machine", RING)]
    tile_w = Inches(2.72)
    gap = Inches(0.28)
    tx0 = MARGIN
    ty = y + Inches(0.08)
    for i, (num, lab, col) in enumerate(stats):
        tx = tx0 + i * (tile_w + gap)
        rect(s, tx, ty, tile_w, Inches(1.95), fill=SURFACE, line=BORDER, line_w=1.25,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, tx + Inches(0.24), ty + Inches(0.26), Inches(0.5), Pt(3.5), fill=col)
        tf = textbox(s, tx + Inches(0.24), ty + Inches(0.46), tile_w - Inches(0.48),
                     Inches(1.4))
        run(para(tf, first=True, after=4), num, size=38, color=col, bold=True,
            font=DISPLAY)
        run(para(tf, after=0, line=1.08), lab, size=12, color=DIM)
    # closing line
    ctf = textbox(s, MARGIN, ty + Inches(2.25), Inches(8.6), Inches(1.3))
    cp = para(ctf, first=True, after=0, line=1.22)
    run(cp, "The people most targeted", size=16, color=WHITE, bold=True)
    run(cp, " — students, retirees, small teams — ", size=16, color=HAZE)
    run(cp, "are the least equipped to judge a link. And the experts who can",
        size=16, color=HAZE)
    run(cp, " are drowning in the queue.", size=16, color=WHITE, bold=True)
    image_fit(s, M_DANGER, SW - Inches(4.25), Inches(3.55), Inches(3.35), Inches(3.35))
    footer(s, n)
    notes(s, """
MICHAEL (~35s). Land the pain, don't recite the numbers. "Phishing is the #1 way breaches start —
one click can detonate. But here's the gap: the people most targeted have the least help, and the
analysts who could help are buried. Today, if you report a suspicious link, you usually hear...
nothing." That last beat sets up the whole demo (the closure loop). Advance.
""")

    # ======================================================================= #
    # 3. HOW ORBIS WORKS (dark)  — Michael
    # ======================================================================= #
    n += 1
    s = add_slide()
    dark_base(s)
    presenter_pill(s, "Michael", MICHAEL)
    y = kicker_head(s, "What Orbis does",
                    "Paste a link. Orbo checks it\nso you never have to.", hsize=29)
    steps = [("1", "Submit", "Paste a URL or forward a suspicious email. No app to install.", RING),
             ("2", "Detonate", "We open it in a secure sandbox and screenshot it — it never touches your machine.", PRIMARY),
             ("3", "Explain", "You get a plain-English verdict and a 0–100 safety score, with the reasons.", SAFE)]
    sy = y + Inches(0.05)
    for i, (num, title, body, col) in enumerate(steps):
        row_y = sy + i * Inches(1.15)
        rect(s, MARGIN, row_y, Inches(0.6), Inches(0.6), fill=None, line=col,
             line_w=2.0, shape=MSO_SHAPE.OVAL)
        ntf = textbox(s, MARGIN, row_y, Inches(0.6), Inches(0.6),
                      anchor=MSO_ANCHOR.MIDDLE)
        run(para(ntf, first=True, align=PP_ALIGN.CENTER, after=0), num, size=20,
            color=col, bold=True, font=DISPLAY)
        ttf = textbox(s, MARGIN + Inches(0.85), row_y - Inches(0.02), Inches(7.4),
                      Inches(1.1))
        run(para(ttf, first=True, after=3), title, size=20, color=WHITE, bold=True,
            font=DISPLAY)
        run(para(ttf, after=0, line=1.1), body, size=13.5, color=DIM)
    image_fit(s, M_THINK, SW - Inches(4.25), Inches(3.35), Inches(3.35), Inches(3.35))
    footer(s, n)
    notes(s, """
MICHAEL (~25s). The mental model for the demo you're about to run — keep it to three words:
"Submit, detonate, explain. You paste a link or forward an email; we detonate it safely in a
sandbox; you get a plain-English verdict and a 0-to-100 score. Let me show you who it's for, then
we'll get into the product." (If running long, you can skip straight past this — it's shown live.)
""")

    # ======================================================================= #
    # 4. THREE EXPERIENCES (dark)  — Michael -> hands to the demo
    # ======================================================================= #
    n += 1
    s = add_slide()
    dark_base(s, right_panel=False)
    y = kicker_head(s, "One codebase, three experiences",
                    "Built for everyone in the blast radius.", hsize=29,
                    width=CW)
    cols = [
        ("Individuals", "Is this link safe?", R_INDIV,
         ["One-tap link & email check", "Plain-language verdict", "Personal history"]),
        ("Org members", "Safe at work, escalate when it matters", R_MEMBR,
         ["Everything individuals get", "Auto-escalate to your analysts", "Shared team history"]),
        ("Analysts", "Expert triage, org-wide", R_ANALY,
         ["Triage queue & campaign view", "Authoritative review that\noverrides the AI", "Ask the data in plain English"]),
    ]
    card_w = (CW - Inches(0.6)) / 3
    ch = Inches(3.85)
    for i, (title, sub, col, items) in enumerate(cols):
        cx = MARGIN + i * (card_w + Inches(0.3))
        rect(s, cx, y, card_w, ch, fill=SURFACE, line=BORDER, line_w=1.25,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, cx + Inches(0.28), y + Inches(0.34), Inches(0.9), Pt(4), fill=col)
        ttf = textbox(s, cx + Inches(0.28), y + Inches(0.55), card_w - Inches(0.56),
                      Inches(0.55))
        run(para(ttf, first=True, after=0), title, size=21, color=WHITE, bold=True,
            font=DISPLAY)
        stf = textbox(s, cx + Inches(0.28), y + Inches(1.08), card_w - Inches(0.56),
                      Inches(0.7))
        run(para(stf, first=True, after=0, line=1.08), sub, size=13, color=col,
            bold=True)
        itf = textbox(s, cx + Inches(0.28), y + Inches(1.95), card_w - Inches(0.5),
                      Inches(1.8))
        for j, it in enumerate(items):
            p = para(itf, first=(j == 0), before=(0 if j == 0 else 9), after=0, line=1.06)
            run(p, "●  ", size=11, color=col)
            run(p, it, size=13, color=DIM)
    footer(s, n)
    notes(s, """
MICHAEL (~30s), then HAND OFF to the demo. "One codebase serves three people: an individual just
wants 'is this safe?'; an org member wants that plus a safety net — their report escalates to their
analyst; and the analyst gets the cockpit — a triage queue, campaign clustering, an authoritative
review that overrides the AI, and they can ask the threat data questions in plain English. Rather
than walk slides, let's show you all of it — David's going to start in the live app."
=> Bring up the browser now. Next slide is the DEMO divider (leave it up or move straight to app).
""")

    # ======================================================================= #
    # 5. LIVE DEMO divider (hero)
    # ======================================================================= #
    n += 1
    s = add_slide()
    hero_bg(s)
    tf = textbox(s, Inches(0.95), Inches(2.0), Inches(8.0), Inches(0.35))
    run(para(tf, first=True, after=0), "LIVE DEMO", size=14, color=ACCENT, bold=True,
        spacing=3.0)
    tf2 = textbox(s, Inches(0.92), Inches(2.5), Inches(8.2), Inches(1.6))
    run(para(tf2, first=True, after=0, line=1.02),
        "Let's show you\nthe real thing.", size=48, color=WHITE, bold=True, font=DISPLAY)
    rect(s, Inches(0.97), Inches(4.35), Inches(2.0), Pt(3), fill=RING)
    tf3 = textbox(s, Inches(0.97), Inches(4.7), Inches(7.6), Inches(0.9))
    p3 = para(tf3, first=True, after=0, line=1.2)
    run(p3, "Submit  ·  read the verdict  ·  escalate  ·  analyst review  ·  ask the data",
        size=16, color=HAZE)
    tf4 = textbox(s, Inches(0.97), Inches(5.6), Inches(7.8), Inches(0.5))
    run(para(tf4, first=True, after=0), "orbis-client-8yu1.onrender.com", size=15,
        color=ACCENT, bold=True)
    image_fit(s, M_WAVE, SW - Inches(4.6), Inches(3.0), Inches(3.9), Inches(3.9))
    footer(s, n)
    notes(s, """
THE DEMO. This divider is a placeholder — you're in the browser now, not on slides. Everything is
PRE-LOADED (no live scans, no waiting). Turn order (each ~even, ~2:20 apart):

DAVID  -> the check-link flagship. Open an already-scanned DANGEROUS result: the verdict card, the
          0-100 score + SAFE/REVIEW/DANGER band, the screenshot, the plain-English "why". Then open
          a SAFE one to contrast. (Optional: the forwarded-email result, to show email triage.)

OZIAS  -> the closure loop. As an org member, open Reports/history and show a submission sitting at
          'pending review' (it auto-escalated). Switch to the analyst account: open the triage queue,
          open that item, submit an authoritative review (score + status 'confirmed malicious' +
          share-with-org). Switch back: the member now sees the confirmed verdict + a notification
          in the bell. THIS is 'from silence to an answer'.

MICHAEL -> scale + the 2nd AI feature. Show the analyst dashboard (stat tiles, trend, donut, campaign
           grouping) reading org-wide data, then Ask-the-data: type a plain-English question on the
           dashboard rail and read the answer/cards it returns.

Keep it PRE-LOADED. Do NOT click 'Sign in with SSO' (not configured). Do NOT submit a brand-new URL.
Then advance to the 'under the hood' slide for the close.
""")

    # ======================================================================= #
    # 6. UNDER THE HOOD (dark)  — David leads the concept, Michael on data layer
    # ======================================================================= #
    n += 1
    s = add_slide()
    dark_base(s, right_panel=False)
    y = kicker_head(s, "Under the hood",
                    "A deterministic core, with an AI narrator.", hsize=29, width=CW)
    # pipeline nodes
    nodes = [("User", "URL or\nforwarded email", RING),
             ("Sandbox", "urlscan.io +\nSafe Browsing", PRIMARY),
             ("Scoring", "our rubric,\n0–100 score", SAFE),
             ("AI", "explains it\nin plain words", REVIEW),
             ("You", "verdict card +\nnext steps", RING)]
    node_w = Inches(2.02)
    node_h = Inches(1.35)
    ngap = (CW - node_w * 5) / 4
    ny = y + Inches(0.05)
    for i, (t, b, col) in enumerate(nodes):
        nx = MARGIN + i * (node_w + ngap)
        rect(s, nx, ny, node_w, node_h, fill=SURFACE, line=col, line_w=1.5,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = textbox(s, nx + Inches(0.14), ny + Inches(0.18), node_w - Inches(0.28),
                     node_h - Inches(0.3))
        run(para(tf, first=True, after=3, align=PP_ALIGN.CENTER), t, size=16,
            color=WHITE, bold=True, font=DISPLAY)
        run(para(tf, after=0, align=PP_ALIGN.CENTER, line=1.08), b, size=11, color=DIM)
        if i < 4:
            atf = textbox(s, nx + node_w, ny, ngap, node_h, anchor=MSO_ANCHOR.MIDDLE)
            run(para(atf, first=True, align=PP_ALIGN.CENTER, after=0), "→", size=20,
                color=ACCENT, bold=True)
    # the differentiator line
    ky = ny + node_h + Inches(0.4)
    rect(s, MARGIN, ky, CW, Inches(0.66), fill=SURFACE, line=RING, line_w=1.5,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    ktf = textbox(s, MARGIN + Inches(0.3), ky, CW - Inches(0.6), Inches(0.66),
                  anchor=MSO_ANCHOR.MIDDLE)
    kp = para(ktf, first=True, after=0)
    run(kp, "Our code owns the score. ", size=15, color=RING, bold=True)
    run(kp, "The AI explains it in plain words — it never overrides the math. "
            "A known-bad link can never read “safe.”", size=15, color=TEXT)
    # two data layers
    ly = ky + Inches(0.95)
    layers = [("Global indicator layer", RING,
               "Each unique URL is scanned once, then shared — fast and cheap on external APIs."),
              ("Per-org data layer", SAFE,
               "Submissions, reviews & campaigns are isolated by organization — a privacy gate by default.")]
    lw = (CW - Inches(0.4)) / 2
    for i, (t, col, b) in enumerate(layers):
        lx = MARGIN + i * (lw + Inches(0.4))
        rect(s, lx, ly, lw, Inches(1.1), fill=SURFACE, line=BORDER, line_w=1.25,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, lx, ly, Inches(0.09), Inches(1.1), fill=col,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = textbox(s, lx + Inches(0.28), ly + Inches(0.16), lw - Inches(0.5),
                     Inches(0.85))
        run(para(tf, first=True, after=4), t, size=15, color=col, bold=True)
        run(para(tf, after=0, line=1.12), b, size=12.5, color=DIM)
    footer(s, n)
    notes(s, """
DAVID leads (~35s); MICHAEL adds the data-layer line. This is the ONE thing you can't show by
clicking, and it's your best answer to "is this just ChatGPT?": "The number is ours — a
deterministic rubric in our code. The AI only writes the explanation; it can nudge the wording but
it can NEVER move a known-bad link into 'safe'. That matters for a security tool." Then Michael:
"And two data layers make it scale and stay private — every unique URL is scanned once and shared,
but each org only ever sees its own submissions and reviews." Advance to the roadmap.
""")

    # ======================================================================= #
    # 7. WHAT'S NEXT (dark, CORRECTED)  — Ozias
    # ======================================================================= #
    n += 1
    s = add_slide()
    dark_base(s, right_panel=False)
    presenter_pill(s, "Ozias", OZIAS)
    y = kicker_head(s, "What's next",
                    "Most of the vision already ships today.", hsize=29, width=CW)
    cols = [
        ("SHIPPING NOW", SAFE, [
            "Link + email triage",
            "3 role experiences",
            "Analyst review that overrides the AI",
            "Notifications — the closure loop",
            "Ask-the-data charts",
            "Campaign clustering in triage",
            "Browser extension (Gmail scan)",
        ]),
        ("NEXT", RING, [
            "Publish the extension to\nthe Chrome Web Store",
            "Real-time alerting\n(today: in-app + emailed reports)",
            "Auto-cluster new submissions\ninto campaigns",
            "Enterprise SSO / SAML",
        ]),
        ("LATER", PRIMARY, [
            "Configurable auto-escalation rules",
            "Public API for SOC / SIEM tools",
            "Email / SMS alerting",
        ]),
    ]
    lane_w = (CW - Inches(0.5)) / 3
    lane_h = Inches(4.05)
    for i, (title, col, items) in enumerate(cols):
        lx = MARGIN + i * (lane_w + Inches(0.25))
        rect(s, lx, y, lane_w, lane_h, fill=SURFACE, line=BORDER, line_w=1.25,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, lx, y, lane_w, Inches(0.52), fill=col,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, lx, y + Inches(0.27), lane_w, Inches(0.25), fill=col)
        htf = textbox(s, lx, y, lane_w, Inches(0.52), anchor=MSO_ANCHOR.MIDDLE)
        run(para(htf, first=True, align=PP_ALIGN.CENTER, after=0), title, size=13,
            color=CANVAS, bold=True, spacing=1.4)
        btf = textbox(s, lx + Inches(0.26), y + Inches(0.72), lane_w - Inches(0.5),
                      lane_h - Inches(0.9))
        for j, it in enumerate(items):
            p = para(btf, first=(j == 0), before=(0 if j == 0 else 8), after=0, line=1.05)
            run(p, "●  ", size=10, color=col)
            run(p, it, size=12.5, color=TEXT)
    footer(s, n)
    notes(s, """
OZIAS (~25s). The honest headline: we've already shipped most of what a roadmap would normally
'promise'. "Everything in the green column is live today — link AND email triage, all three roles,
the analyst review that overrides the AI, notifications, ask-the-data, campaign clustering, and a
browser extension that scans Gmail inline. What's genuinely next is polish and reach: publishing the
extension, real-time push alerts (today it's in-app plus emailed reports), and enterprise SSO."
Be honest — say 'in-app + emailed reports' so 'real-time' reads as an upgrade, not a fix. Advance.
""")

    # ======================================================================= #
    # 8. CONCLUSION (hero)  — Michael closes
    # ======================================================================= #
    n += 1
    s = add_slide()
    hero_bg(s)
    tf = textbox(s, Inches(0.95), Inches(0.85), Inches(8.0), Inches(0.35))
    run(para(tf, first=True, after=0), "IN SHORT", size=14, color=ACCENT, bold=True,
        spacing=3.0)
    tf2 = textbox(s, Inches(0.92), Inches(1.35), Inches(8.2), Inches(1.4))
    run(para(tf2, first=True, after=0, line=1.03),
        "From silence to a verdict —\nfor everyone in the blast radius.", size=33,
        color=WHITE, bold=True, font=DISPLAY)
    rect(s, Inches(0.97), Inches(2.95), Inches(2.0), Pt(3), fill=RING)
    pairs = [
        ("The problem", "Phishing is a human problem — but safe triage lived only in expert tools, and reporters never heard back."),
        ("What we built", "Paste a link or forward an email → a sandboxed, deterministic verdict in plain English, for anyone."),
        ("The closure loop", "An analyst's authoritative review overrides the AI and notifies the person who reported it."),
        ("Why it holds up", "Our code owns the score; the AI only narrates. Shared intel, per-org privacy, by default."),
    ]
    py = Inches(3.4)
    for i, (h, b) in enumerate(pairs):
        rect(s, Inches(0.97), py + Inches(0.05), Inches(0.14), Inches(0.14),
             fill=RING, shape=MSO_SHAPE.OVAL)
        tf = textbox(s, Inches(1.35), py - Inches(0.06), Inches(7.5), Inches(0.85))
        p = para(tf, first=True, after=0, line=1.08)
        run(p, h + "   ", size=15, color=WHITE, bold=True)
        run(p, b, size=14, color=HAZE)
        py += Inches(0.82)
    image_fit(s, M_HAPPY, SW - Inches(4.7), Inches(2.7), Inches(4.0), Inches(4.0))
    footer(s, n)
    notes(s, """
MICHAEL closes (~35s). Tie the bow back to the problem you opened with: "We started with silence —
you report a link and hear nothing. Orbis turns that into an answer: paste or forward, get a
plain-English verdict anyone can act on, and when an analyst confirms it, the person who reported it
finally gets told. Expert-grade triage, made one-tap simple." Pause, then advance to Thank you.
""")

    # ======================================================================= #
    # 9. THANK YOU (hero)
    # ======================================================================= #
    n += 1
    s = add_slide()
    hero_bg(s)
    image_fit(s, LOGO_DARK, Inches(0.9), Inches(0.9), Inches(3.1), Inches(0.95))
    tf = textbox(s, Inches(0.92), Inches(2.7), Inches(8.0), Inches(1.2))
    run(para(tf, first=True, after=0), "Thank you.", size=60, color=WHITE, bold=True,
        font=DISPLAY)
    rect(s, Inches(0.97), Inches(4.05), Inches(2.0), Pt(3), fill=RING)
    tf2 = textbox(s, Inches(0.97), Inches(4.35), Inches(7.6), Inches(1.0))
    run(para(tf2, first=True, after=0, line=1.2),
        "Expert-grade phishing triage, made accessible.\nWe'd love your questions.",
        size=17, color=HAZE)
    tf3 = textbox(s, Inches(0.97), Inches(5.75), Inches(9.2), Inches(0.5))
    p3 = para(tf3, first=True, after=0)
    run(p3, "Team DOMinion", size=13, color=WHITE, bold=True)
    run(p3, "   ·   Michael Jissa  ·  Ozias Tumimana  ·  David Gonzalez-Cesar",
        size=12.5, color=HAZE)
    tf4 = textbox(s, Inches(0.97), Inches(6.3), Inches(8), Inches(0.4))
    run(para(tf4, first=True, after=0), "orbis-client-8yu1.onrender.com", size=14,
        color=ACCENT, bold=True)
    image_fit(s, M_SAFE, SW - Inches(4.55), Inches(3.1), Inches(3.7), Inches(3.7))
    footer(s, n)
    notes(s, """
ALL. Open the floor. Keep the app up so you can answer "show me X" live. If they poke at the
roadmap: SSO button exists but isn't configured (don't click it), notifications are in-app +
emailed reports today (real-time push is next), and campaign clustering groups from seeded data
(auto-clustering new submissions is next). One of you: capture their questions as action items.
""")

    prs.save(OUT)
    count = len(prs.slides._sldIdLst)
    print(f"Wrote {OUT} with {count} slides")


if __name__ == "__main__":
    build()
